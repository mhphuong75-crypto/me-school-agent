#!/usr/bin/env python3
"""
ME School Manual — Document Ingestion Script

Two modes (auto-detected from .env):
  • SharePoint mode  — reads files directly from SharePoint via Microsoft Graph API
                       (used when SHAREPOINT_TENANT_ID etc. are set in .env)
  • Local mode       — reads files from MANUAL_ROOT folder on this machine

Usage:
    python ingest.py
    python ingest.py --reset   # wipe and rebuild from scratch
"""

import os, sys, tempfile, json, time
from pathlib import Path
from dotenv import load_dotenv

load_dotenv()

# ── Anthropic (for TOC generation) ────────────────────────────────────────
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY", "")
TOC_MODEL = "claude-haiku-4-5"

# ── Config ─────────────────────────────────────────────────────────────────

# Local mode
MANUAL_ROOT   = Path(os.getenv("MANUAL_ROOT", "./OneDrive_1_10-05-2026"))
ONEDRIVE_BASE = os.getenv("ONEDRIVE_BASE_URL", "").rstrip("/")

# SharePoint mode (all four must be set to activate SharePoint mode)
SP_TENANT_ID     = os.getenv("SHAREPOINT_TENANT_ID", "")
SP_CLIENT_ID     = os.getenv("SHAREPOINT_CLIENT_ID", "")
SP_CLIENT_SECRET = os.getenv("SHAREPOINT_CLIENT_SECRET", "")
SP_SITE_URL      = os.getenv("SHAREPOINT_SITE_URL", "")
SP_FOLDER        = os.getenv("SHAREPOINT_FOLDER", "")   # relative path inside document library
                                                         # e.g. "Sổ tay/ME School Manual_Sotana"

USE_SHAREPOINT = all([SP_TENANT_ID, SP_CLIENT_ID, SP_CLIENT_SECRET, SP_SITE_URL])

CHUNK_WORDS   = 600
CHUNK_OVERLAP = 80
SUPPORTED_EXT = {".docx", ".doc", ".pdf", ".xlsx", ".xls", ".pptx", ".ppt", ".txt"}


# ── Text extractors ────────────────────────────────────────────────────────

def extract_docx(path: Path) -> str:
    import docx
    try:
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        print(f"  WARN docx {path.name}: {e}")
        return ""

def extract_pdf(path: Path) -> str:
    import pdfplumber
    try:
        with pdfplumber.open(str(path)) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)
    except Exception as e:
        print(f"  WARN pdf {path.name}: {e}")
        return ""

def extract_xlsx(path: Path) -> str:
    import openpyxl
    try:
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        rows = []
        for ws in wb.worksheets:
            rows.append(f"[Sheet: {ws.title}]")
            for row in ws.iter_rows():
                vals = [str(c.value) for c in row if c.value is not None]
                if vals:
                    rows.append(" | ".join(vals))
        return "\n".join(rows)
    except Exception as e:
        print(f"  WARN xlsx {path.name}: {e}")
        return ""

def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    try:
        prs = Presentation(str(path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        print(f"  WARN pptx {path.name}: {e}")
        return ""

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in (".docx", ".doc"):   return extract_docx(path)
    elif ext == ".pdf":             return extract_pdf(path)
    elif ext in (".xlsx", ".xls"): return extract_xlsx(path)
    elif ext in (".pptx", ".ppt"): return extract_pptx(path)
    elif ext == ".txt":            return path.read_text(errors="ignore")
    return ""


# ── Chunking ───────────────────────────────────────────────────────────────

def chunk(text: str, size=CHUNK_WORDS, overlap=CHUNK_OVERLAP):
    words = text.split()
    chunks, i = [], 0
    while i < len(words):
        chunks.append(" ".join(words[i:i + size]))
        i += size - overlap
    return [c for c in chunks if len(c.strip()) > 30]


# ── TOC generation (LLM-powered) ──────────────────────────────────────────

def _get_anthropic_client():
    """Lazy init Anthropic client."""
    import anthropic
    return anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

def generate_toc_entry(text: str, file_name: str, file_path: str) -> dict | None:
    """
    Call Claude Haiku to generate a summary + answerable questions for one file.
    Returns dict with summary + questions, or None on failure.
    """
    if not ANTHROPIC_API_KEY:
        return None

    # Truncate very long texts to ~8000 words to stay within Haiku context
    words = text.split()
    truncated = " ".join(words[:8000]) if len(words) > 8000 else text

    prompt = f"""Bạn là trợ lý phân tích tài liệu vận hành của ME School (trường mầm non).

FILE: {file_name}
PATH: {file_path}

NỘI DUNG:
{truncated}

NHIỆM VỤ: Phân tích tài liệu trên và trả về JSON (CHỈ JSON, không text khác):
{{
  "summary": "Tóm tắt nội dung chính của tài liệu trong 2-3 câu bằng tiếng Việt",
  "questions": ["Liệt kê 5-10 câu hỏi mà tài liệu này CÓ THỂ trả lời được, viết bằng tiếng Việt, dùng ngôn ngữ tự nhiên như cách nhân viên sẽ hỏi"]
}}

YÊU CẦU cho questions:
- Viết như cách nhân viên thật sẽ hỏi (ví dụ: "Quy trình xin nghỉ phép như thế nào?")
- Bao gồm cả câu hỏi ngắn lẫn câu hỏi chi tiết
- Nếu tài liệu có nhiều chủ đề, mỗi chủ đề ít nhất 1 câu hỏi
- Dùng từ khóa quan trọng trong tài liệu"""

    try:
        client = _get_anthropic_client()
        resp = client.messages.create(
            model=TOC_MODEL,
            max_tokens=1000,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = resp.content[0].text.strip()
        # Clean markdown code blocks if present
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        if raw.endswith("```"):
            raw = raw[:-3]
        data = json.loads(raw.strip())
        return {
            "summary": data.get("summary", ""),
            "questions": data.get("questions", []),
        }
    except Exception as e:
        print(f"  WARN TOC generation failed for {file_name}: {type(e).__name__}: {e}")
        return None


# ── SharePoint / Microsoft Graph helpers ──────────────────────────────────

def _graph_token() -> str:
    """Get Microsoft Graph access token via client credentials."""
    from msal import ConfidentialClientApplication
    app = ConfidentialClientApplication(
        SP_CLIENT_ID,
        authority=f"https://login.microsoftonline.com/{SP_TENANT_ID}",
        client_credential=SP_CLIENT_SECRET,
    )
    result = app.acquire_token_for_client(
        scopes=["https://graph.microsoft.com/.default"]
    )
    if "access_token" not in result:
        raise RuntimeError(
            f"Microsoft authentication failed: {result.get('error_description', result)}"
        )
    return result["access_token"]


def _graph_get(token: str, url: str) -> dict:
    import requests
    resp = requests.get(url, headers={"Authorization": f"Bearer {token}"})
    resp.raise_for_status()
    return resp.json()


def _get_drive_id(token: str) -> str:
    """Get the default document library drive ID for the configured SharePoint site."""
    from urllib.parse import urlparse
    parsed   = urlparse(SP_SITE_URL)
    hostname = parsed.netloc
    path     = parsed.path.rstrip("/")
    site     = _graph_get(token, f"https://graph.microsoft.com/v1.0/sites/{hostname}:{path}")
    drive    = _graph_get(token, f"https://graph.microsoft.com/v1.0/sites/{site['id']}/drive")
    return drive["id"]


def _list_files(token: str, drive_id: str, folder: str) -> list:
    """Recursively list all supported files under folder path."""
    if folder:
        encoded = folder.replace(" ", "%20")
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root:/{encoded}:/children"
    else:
        url = f"https://graph.microsoft.com/v1.0/drives/{drive_id}/root/children"

    files = []
    while url:
        data = _graph_get(token, url)
        for item in data.get("value", []):
            item_path = f"{folder}/{item['name']}" if folder else item["name"]
            if "folder" in item:
                files.extend(_list_files(token, drive_id, item_path))
            elif "file" in item:
                if Path(item["name"]).suffix.lower() in SUPPORTED_EXT:
                    files.append({
                        "name":         item["name"],
                        "path":         item_path,
                        "download_url": item.get("@microsoft.graph.downloadUrl", ""),
                        "web_url":      item.get("webUrl", ""),
                    })
        url = data.get("@odata.nextLink")   # pagination
    return files


def _download(url: str, dest: Path):
    import requests
    resp = requests.get(url, stream=True)
    resp.raise_for_status()
    with open(dest, "wb") as f:
        for block in resp.iter_content(chunk_size=8192):
            f.write(block)


# ── Main — SharePoint mode ─────────────────────────────────────────────────

def main_sharepoint():
    import numpy as np
    from sentence_transformers import SentenceTransformer

    print("Mode     : SharePoint (Microsoft Graph API)")
    print(f"Site     : {SP_SITE_URL}")
    print(f"Folder   : {SP_FOLDER or '(root of document library)'}")
    print(f"TOC      : {'enabled (Haiku)' if ANTHROPIC_API_KEY else 'disabled (no API key)'}")
    print()

    print("Authenticating with Microsoft…")
    token = _graph_token()
    print("✓ Authenticated\n")

    print("Connecting to SharePoint…")
    drive_id = _get_drive_id(token)
    print("✓ Connected\n")

    print("Listing files…")
    files = _list_files(token, drive_id, SP_FOLDER)
    print(f"Found {len(files)} supported files.\n")

    if not files:
        print("ERROR: No files found. Check SHAREPOINT_FOLDER path in .env")
        sys.exit(1)

    print("Loading embedding model (first run downloads ~120 MB)…")
    model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")
    print("Model ready.\n")

    all_vectors, all_records = [], []
    toc_entries = []
    toc_success, toc_fail = 0, 0

    with tempfile.TemporaryDirectory() as tmpdir:
        for idx, f in enumerate(files, 1):
            print(f"[{idx}/{len(files)}] {f['path']}")

            tmp_path = Path(tmpdir) / f["name"]
            try:
                _download(f["download_url"], tmp_path)
            except Exception as e:
                print(f"  WARN download failed: {e}")
                continue

            text = extract_text(tmp_path)
            if not text.strip():
                print("  → (no text, skipping)")
                continue

            chunks = chunk(text)
            if not chunks:
                continue

            folder_label = f["path"].replace("\\", "/").split("/")[0]
            embeddings   = model.encode(chunks, show_progress_bar=False)

            for c, emb in zip(chunks, embeddings):
                all_vectors.append(emb)
                all_records.append({
                    "text":      c,
                    "file_name": f["name"],
                    "file_path": f["path"],
                    "folder":    folder_label,
                    "url":       f["web_url"],
                })

            # ── TOC generation (with fail-fast) ──
            toc_enabled = ANTHROPIC_API_KEY and not (toc_fail >= 3 and toc_success == 0)
            if toc_enabled:
                toc = generate_toc_entry(text, f["name"], f["path"])
                if toc:
                    toc_entries.append({
                        "file_name": f["name"],
                        "file_path": f["path"],
                        "folder":    folder_label,
                        "url":       f["web_url"],
                        "summary":   toc["summary"],
                        "questions": toc["questions"],
                        "word_count": len(text.split()),
                    })
                    toc_success += 1
                    print(f"  → {len(chunks)} chunks + TOC ✓")
                else:
                    toc_fail += 1
                    if toc_fail >= 3 and toc_success == 0:
                        print(f"  → {len(chunks)} chunks (TOC failed)")
                        print(f"  ⚠ TOC DISABLED: first {toc_fail} calls all failed — skipping TOC for remaining files")
                    else:
                        print(f"  → {len(chunks)} chunks (TOC failed)")
                # Rate limit: small delay between Haiku calls
                time.sleep(0.3)
            else:
                print(f"  → {len(chunks)} chunks indexed")

    if ANTHROPIC_API_KEY:
        print(f"\nTOC: {toc_success} success, {toc_fail} failed out of {toc_success + toc_fail} files")
    _save(all_vectors, all_records, toc_entries)


# ── Main — Local mode ──────────────────────────────────────────────────────

def main_local():
    import numpy as np
    from sentence_transformers import SentenceTransformer

    print("Mode     : Local folder")
    print(f"Root     : {MANUAL_ROOT}")
    print(f"URL base : {ONEDRIVE_BASE or '(not set — file:// links used)'}")
    print(f"TOC      : {'enabled (Haiku)' if ANTHROPIC_API_KEY else 'disabled (no API key)'}")
    print()

    if not MANUAL_ROOT.exists():
        print(f"ERROR: MANUAL_ROOT not found: {MANUAL_ROOT}")
        sys.exit(1)

    files = [
        f for f in MANUAL_ROOT.rglob("*")
        if f.is_file()
        and f.suffix.lower() in SUPPORTED_EXT
        and not f.name.startswith((".", "~$", "__TEMP"))
        and f.name != ".DS_Store"
    ]
    print(f"Found {len(files)} files.\n")

    print("Loading embedding model (first run downloads ~120 MB)…")
    model = SentenceTransformer("paraphrase-multilingual-MiniLM-L12-v2")
    print("Model ready.\n")

    all_vectors, all_records = [], []
    toc_entries = []
    toc_success, toc_fail = 0, 0

    for idx, fpath in enumerate(files, 1):
        rel    = fpath.relative_to(MANUAL_ROOT)
        folder = rel.parts[0] if len(rel.parts) > 1 else "root"
        print(f"[{idx}/{len(files)}] {rel}")

        text = extract_text(fpath)
        if not text.strip():
            print("  → (no text, skipping)")
            continue

        chunks = chunk(text)
        if not chunks:
            continue

        encoded    = str(rel).replace("\\", "/").replace(" ", "%20")
        url        = f"{ONEDRIVE_BASE}/{encoded}" if ONEDRIVE_BASE else f"file://{MANUAL_ROOT / rel}"
        embeddings = model.encode(chunks, show_progress_bar=False)

        for c, emb in zip(chunks, embeddings):
            all_vectors.append(emb)
            all_records.append({
                "text":      c,
                "file_name": fpath.name,
                "file_path": str(rel),
                "folder":    folder,
                "url":       url,
            })

        # ── TOC generation ──
        if ANTHROPIC_API_KEY:
            toc = generate_toc_entry(text, fpath.name, str(rel))
            if toc:
                toc_entries.append({
                    "file_name": fpath.name,
                    "file_path": str(rel),
                    "folder":    folder,
                    "url":       url,
                    "summary":   toc["summary"],
                    "questions": toc["questions"],
                    "word_count": len(text.split()),
                })
                toc_success += 1
                print(f"  → {len(chunks)} chunks + TOC ✓")
            else:
                toc_fail += 1
                print(f"  → {len(chunks)} chunks (TOC failed)")
            time.sleep(0.3)
        else:
            print(f"  → {len(chunks)} chunks indexed")

    if ANTHROPIC_API_KEY:
        print(f"\nTOC: {toc_success} success, {toc_fail} failed out of {toc_success + toc_fail} files")
    _save(all_vectors, all_records, toc_entries)


# ── Save ───────────────────────────────────────────────────────────────────

def _save(all_vectors: list, all_records: list, toc_entries: list | None = None):
    import numpy as np

    vectors_path  = Path("vectors.npy")
    metadata_path = Path("metadata.json")
    toc_path      = Path("toc.json")

    np.save(str(vectors_path), np.array(all_vectors, dtype="float32"))
    with open(metadata_path, "w", encoding="utf-8") as f:
        json.dump(all_records, f, ensure_ascii=False)

    print(f"\n✓ Done.  Total chunks : {len(all_records)}")
    print(f"  Saved  : {vectors_path}  ({vectors_path.stat().st_size/1024/1024:.1f} MB)")
    print(f"  Saved  : {metadata_path} ({metadata_path.stat().st_size/1024/1024:.1f} MB)")

    if toc_entries:
        with open(toc_path, "w", encoding="utf-8") as f:
            json.dump(toc_entries, f, ensure_ascii=False, indent=2)
        print(f"  Saved  : {toc_path}  ({toc_path.stat().st_size/1024:.0f} KB, {len(toc_entries)} files)")
    else:
        print("  TOC    : not generated (no API key or all failed)")


# ── Entry point ────────────────────────────────────────────────────────────

if __name__ == "__main__":
    if USE_SHAREPOINT:
        main_sharepoint()
    else:
        main_local()
